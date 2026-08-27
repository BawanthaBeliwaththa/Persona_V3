import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette
    COLOR_BG = RGBColor(11, 15, 25)        # Dark Slate #0b0f19
    COLOR_CARD = RGBColor(20, 29, 47)      # Deep Card #141d2f
    COLOR_BLUE = RGBColor(10, 102, 194)    # LinkedIn Blue #0a66c2
    COLOR_PURPLE = RGBColor(99, 102, 241)  # Indigo #6366f1
    COLOR_WHITE = RGBColor(243, 244, 246)  # White #f3f4f6
    COLOR_MUTED = RGBColor(156, 163, 175)  # Muted Gray #9ca3af
    COLOR_GREEN = RGBColor(16, 185, 129)   # Emerald #10b981
    COLOR_ACCENT = RGBColor(96, 165, 250)  # Light Blue #60a5fa

    blank_layout = prs.slide_layouts[6]

    def add_header(slide, title_text, category_text="PERSONA V3 GUIDE"):
        # Background rect
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_BG
        bg.line.fill.background()

        # Category badge
        badge = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        tf_b = badge.text_frame
        p_b = tf_b.paragraphs[0]
        p_b.text = category_text.upper()
        p_b.font.size = Pt(11)
        p_b.font.bold = True
        p_b.font.color.rgb = COLOR_ACCENT

        # Title
        title = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(11.7), Inches(0.8))
        tf_t = title.text_frame
        p_t = tf_t.paragraphs[0]
        p_t.text = title_text
        p_t.font.size = Pt(24)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_WHITE

    # ----------------------------------------------------
    # SLIDE 1: Title Slide
    # ----------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = COLOR_BG
    bg1.line.fill.background()

    card1 = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.2), Inches(11.333), Inches(5.1))
    card1.fill.solid()
    card1.fill.fore_color.rgb = COLOR_CARD
    card1.line.color.rgb = COLOR_BLUE

    tf1 = card1.text_frame
    tf1.word_wrap = True

    p = tf1.paragraphs[0]
    p.text = "PERSONA V3"
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = COLOR_BLUE

    p2 = tf1.add_paragraph()
    p2.text = "Complete System Architecture, User Manual & VPS Deployment Blueprint"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_WHITE
    p2.space_before = Pt(14)

    p3 = tf1.add_paragraph()
    p3.text = "GitHub Repository: https://github.com/BawanthaBeliwaththa/Persona_V3"
    p3.font.size = Pt(14)
    p3.font.color.rgb = COLOR_ACCENT
    p3.space_before = Pt(20)

    p4 = tf1.add_paragraph()
    p4.text = "Includes step-by-step instructions for Local Setup, Portal Operations, Task Bucket Queue, GUI & SSH VPS Deployment, and API Integration."
    p4.font.size = Pt(13)
    p4.font.color.rgb = COLOR_MUTED
    p4.space_before = Pt(14)

    # ----------------------------------------------------
    # SLIDE 2: System Architecture
    # ----------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    add_header(slide2, "1. Persona V3 System Architecture & Tech Stack")

    features = [
        ("Flask & Async Loop Core", "Multi-threaded Flask backend coupled with a background asyncio event loop to handle Playwright Chromium operations without blocking HTTP requests."),
        ("Playwright Chromium Engine", "High-performance headless browser automation for scraping profile details, experience, skills, certifications, and recommendation text."),
        ("Task Bucket Queue", "Automated rate-limited queue processing system. Prevents LinkedIn account blocks by applying configurable rest timers (default: 30s) between scrapes."),
        ("Real-Time SSE Event Stream", "Server-Sent Events (SSE) push live terminal logs, worker rest countdowns, and queue status updates directly to admin and client dashboards.")
    ]

    for idx, (title, desc) in enumerate(features):
        row = idx // 2
        col = idx % 2
        x = Inches(0.8 + col * 5.9)
        y = Inches(1.8 + row * 2.6)

        c = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.6), Inches(2.3))
        c.fill.solid()
        c.fill.fore_color.rgb = COLOR_CARD
        c.line.color.rgb = COLOR_PURPLE

        tf = c.text_frame
        tf.word_wrap = True
        p_head = tf.paragraphs[0]
        p_head.text = title
        p_head.font.size = Pt(16)
        p_head.font.bold = True
        p_head.font.color.rgb = COLOR_WHITE

        p_body = tf.add_paragraph()
        p_body.text = desc
        p_body.font.size = Pt(12)
        p_body.font.color.rgb = COLOR_MUTED
        p_body.space_before = Pt(8)

    # ----------------------------------------------------
    # SLIDE 3: How to Install & Run Locally
    # ----------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    add_header(slide3, "2. How to Install and Run Persona V3 Locally")

    steps = [
        ("Step 1: Clone Repository", "git clone https://github.com/BawanthaBeliwaththa/Persona_V3.git\ncd Persona_V3"),
        ("Step 2: Create Virtual Environment", "python -m venv venv\nvenv\\Scripts\\activate  (or source venv/bin/activate on Linux)"),
        ("Step 3: Install Dependencies", "pip install --upgrade pip\npip install -r req.txt"),
        ("Step 4: Install Chromium Browser", "playwright install chromium"),
        ("Step 5: Run Application", "python app.py  --> Access at http://localhost:5000")
    ]

    for idx, (stitle, scode) in enumerate(steps):
        y = Inches(1.6 + idx * 1.05)
        c = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(11.7), Inches(0.95))
        c.fill.solid()
        c.fill.fore_color.rgb = COLOR_CARD
        c.line.color.rgb = COLOR_BLUE

        tf = c.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = stitle
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_ACCENT

        p2 = tf.add_paragraph()
        p2.text = scode
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_WHITE
        p2.space_before = Pt(4)

    # ----------------------------------------------------
    # SLIDE 4: Admin & Client Portals
    # ----------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    add_header(slide4, "3. Operational Guide: Admin & Client Portals")

    # Left: Admin
    c_adm = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0))
    c_adm.fill.solid()
    c_adm.fill.fore_color.rgb = COLOR_CARD
    c_adm.line.color.rgb = COLOR_PURPLE

    tf_a = c_adm.text_frame
    tf_a.word_wrap = True
    pa = tf_a.paragraphs[0]
    pa.text = "ADMIN PORTAL (/admin)"
    pa.font.size = Pt(18)
    pa.font.bold = True
    pa.font.color.rgb = COLOR_WHITE

    items_admin = [
        "Initialize Scraper: Select Headless (On/Off) & browser type.",
        "LinkedIn Login: Authenticate active session credentials.",
        "Task Bucket Control: Pause/Resume worker, set rest timer (sec).",
        "Master Database: Export all non-volatile profile records in JSON & CSV formats.",
        "Force Kill Browser: Emergency process termination button."
    ]
    for it in items_admin:
        p_i = tf_a.add_paragraph()
        p_i.text = "• " + it
        p_i.font.size = Pt(12)
        p_i.font.color.rgb = COLOR_MUTED
        p_i.space_before = Pt(10)

    # Right: Client
    c_cli = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.8), Inches(5.6), Inches(5.0))
    c_cli.fill.solid()
    c_cli.fill.fore_color.rgb = COLOR_CARD
    c_cli.line.color.rgb = COLOR_BLUE

    tf_c = c_cli.text_frame
    tf_c.word_wrap = True
    pc = tf_c.paragraphs[0]
    pc.text = "CLIENT PORTAL (/)"
    pc.font.size = Pt(18)
    pc.font.bold = True
    pc.font.color.rgb = COLOR_WHITE

    items_client = [
        "Direct Profile Extraction: Enter username or full LinkedIn URL.",
        "Reference Number Tracking: Unique REF-XXXXXX code generated per search for easy reload.",
        "Bulk CSV / JSON Tasks: Upload multiple search queries or URLs.",
        "PDF Report Export: Print & export candidate profile reports.",
        "Contact Search: Search by email/phone for premium lookups."
    ]
    for it in items_client:
        p_i = tf_c.add_paragraph()
        p_i.text = "• " + it
        p_i.font.size = Pt(12)
        p_i.font.color.rgb = COLOR_MUTED
        p_i.space_before = Pt(10)

    # ----------------------------------------------------
    # SLIDE 5: Task Bucket Queue
    # ----------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    add_header(slide5, "4. Task Bucket Queue System & Anti-Ban Protection")

    tb_cards = [
        ("Sequential Queue Execution", "All bulk searches and direct URL tasks are added to the Task Bucket. The worker executes tasks one-by-one to avoid rate-limiting."),
        ("Configurable Rest Period", "After completing each profile extraction, the worker pauses for a customizable rest period (default: 30 seconds) to mimic human behavior."),
        ("Worker Controls", "Admins can Pause, Resume, or Clear completed/failed tasks at any time from the Task Bucket control bar."),
        ("Bulk Task Import", "Upload CSV or JSON files containing hundreds of usernames or URLs into the bucket in a single click.")
    ]
    for idx, (title, desc) in enumerate(tb_cards):
        row = idx // 2
        col = idx % 2
        x = Inches(0.8 + col * 5.9)
        y = Inches(1.8 + row * 2.6)

        c = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.6), Inches(2.3))
        c.fill.solid()
        c.fill.fore_color.rgb = COLOR_CARD
        c.line.color.rgb = COLOR_GREEN

        tf = c.text_frame
        tf.word_wrap = True
        p_head = tf.paragraphs[0]
        p_head.text = title
        p_head.font.size = Pt(16)
        p_head.font.bold = True
        p_head.font.color.rgb = COLOR_WHITE

        p_body = tf.add_paragraph()
        p_body.text = desc
        p_body.font.size = Pt(12)
        p_body.font.color.rgb = COLOR_MUTED
        p_body.space_before = Pt(8)

    # ----------------------------------------------------
    # SLIDE 6: VPS Deployment (GUI / Coolify)
    # ----------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    add_header(slide6, "5. VPS Deployment: GUI & PaaS Method (No Terminal)")

    gui_steps = [
        ("1. Create VPS Instance", "Launch an Ubuntu VPS on DigitalOcean, Hetzner, or Vultr using the 1-Click Coolify App option."),
        ("2. Open Web Dashboard", "Access https://YOUR_VPS_IP:8000 in your browser and register your admin account."),
        ("3. Connect GitHub Repo", "Go to Sources -> GitHub App, authorize Coolify, and select repository BawanthaBeliwaththa/Persona_V3."),
        ("4. Auto Container Build", "Coolify detects the repository Dockerfile, runs 'playwright install chromium', provisions HTTPS SSL, and deploys automatic updates on git push.")
    ]
    for idx, (stitle, sdesc) in enumerate(gui_steps):
        y = Inches(1.6 + idx * 1.3)
        c = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(11.7), Inches(1.15))
        c.fill.solid()
        c.fill.fore_color.rgb = COLOR_CARD
        c.line.color.rgb = COLOR_PURPLE

        tf = c.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = stitle
        p1.font.size = Pt(15)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_ACCENT

        p2 = tf.add_paragraph()
        p2.text = sdesc
        p2.font.size = Pt(12)
        p2.font.color.rgb = COLOR_MUTED
        p2.space_before = Pt(4)

    # ----------------------------------------------------
    # SLIDE 7: VPS Deployment (CLI / SSH)
    # ----------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    add_header(slide7, "6. VPS Deployment: Terminal Method (SSH / Nginx)")

    cli_steps = [
        ("Step 1: Install System Packages", "ssh root@YOUR_VPS_IP\nsudo apt update && sudo apt install -y python3 python3-pip python3-venv git nginx"),
        ("Step 2: Clone Repo & Setup Venv", "cd /var/www && git clone https://github.com/BawanthaBeliwaththa/Persona_V3.git\ncd Persona_V3 && python3 -m venv venv && source venv/bin/activate && pip install -r req.txt"),
        ("Step 3: Install Playwright OS Libs", "playwright install --with-deps chromium  (Installs X11 & font libraries for headless execution)"),
        ("Step 4: Configure Systemd & Nginx", "Create /etc/systemd/system/personav3.service & Nginx site config with SSE buffering disabled.\nsudo systemctl start personav3 && sudo certbot --nginx -d api.yourdomain.com")
    ]
    for idx, (stitle, sdesc) in enumerate(cli_steps):
        y = Inches(1.6 + idx * 1.3)
        c = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(11.7), Inches(1.15))
        c.fill.solid()
        c.fill.fore_color.rgb = COLOR_CARD
        c.line.color.rgb = COLOR_BLUE

        tf = c.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = stitle
        p1.font.size = Pt(15)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_ACCENT

        p2 = tf.add_paragraph()
        p2.text = sdesc
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_WHITE
        p2.space_before = Pt(4)

    # ----------------------------------------------------
    # SLIDE 8: Frontend API Integration
    # ----------------------------------------------------
    slide8 = prs.slides.add_slide(blank_layout)
    add_header(slide8, "7. Frontend Team API Integration Architecture")

    # Left: Security & CORS
    c_sec = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0))
    c_sec.fill.solid()
    c_sec.fill.fore_color.rgb = COLOR_CARD
    c_sec.line.color.rgb = COLOR_BLUE

    tf_s = c_sec.text_frame
    tf_s.word_wrap = True
    ps = tf_s.paragraphs[0]
    ps.text = "CORS & SECURITY CONFIGURATION"
    ps.font.size = Pt(16)
    ps.font.bold = True
    ps.font.color.rgb = COLOR_WHITE

    sec_items = [
        "flask_cors initialized in app.py for cross-origin AJAX requests.",
        "Restrict allowed origins in production:",
        "CORS(app, resources={r'/api/*': {'origins': ['https://frontend-app.com']}})",
        "API Authorization Header Middleware:",
        "Authorization: Bearer YOUR_API_SECRET_TOKEN"
    ]
    for it in sec_items:
        p_i = tf_s.add_paragraph()
        p_i.text = "• " + it
        p_i.font.size = Pt(12)
        p_i.font.color.rgb = COLOR_MUTED
        p_i.space_before = Pt(8)

    # Right: Endpoints Table
    c_ep = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.8), Inches(5.6), Inches(5.0))
    c_ep.fill.solid()
    c_ep.fill.fore_color.rgb = COLOR_CARD
    c_ep.line.color.rgb = COLOR_PURPLE

    tf_e = c_ep.text_frame
    tf_e.word_wrap = True
    pe = tf_e.paragraphs[0]
    pe.text = "CORE API ENDPOINTS"
    pe.font.size = Pt(16)
    pe.font.bold = True
    pe.font.color.rgb = COLOR_WHITE

    endpoints = [
        ("POST /api/scraper/init", "Initializes Playwright Chromium session."),
        ("POST /api/client/scrape", "Extracts LinkedIn profile details by name/URL."),
        ("POST /api/scraper/kill-browser", "Force-kills stuck browser processes."),
        ("GET /api/events", "Server-Sent Events (SSE) live progress stream.")
    ]
    for route, desc in endpoints:
        p_r = tf_e.add_paragraph()
        p_r.text = "• " + route
        p_r.font.size = Pt(13)
        p_r.font.bold = True
        p_r.font.color.rgb = COLOR_ACCENT
        p_r.space_before = Pt(8)

        p_d = tf_e.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = COLOR_MUTED

    # ----------------------------------------------------
    # SLIDE 9: Code Integration Examples
    # ----------------------------------------------------
    slide9 = prs.slides.add_slide(blank_layout)
    add_header(slide9, "8. Frontend Integration Code Examples")

    c_code1 = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0))
    c_code1.fill.solid()
    c_code1.fill.fore_color.rgb = COLOR_CARD
    c_code1.line.color.rgb = COLOR_GREEN

    tf_c1 = c_code1.text_frame
    tf_c1.word_wrap = True
    p = tf_c1.paragraphs[0]
    p.text = "1. FETCH PROFILE EXTRACTION (JS)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    code1 = """const API_BASE = "https://api.yourdomain.com";

async function scrapeProfile(urlOrName) {
  const res = await fetch(`${API_BASE}/api/client/scrape`, {
    method: 'POST',
    headers: {
      'Content-Type': 'application/json',
      'Authorization': 'Bearer TOKEN'
    },
    body: JSON.stringify({ name: urlOrName })
  });
  return await res.json();
}"""
    p_code = tf_c1.add_paragraph()
    p_code.text = code1
    p_code.font.size = Pt(11)
    p_code.font.color.rgb = COLOR_ACCENT
    p_code.space_before = Pt(10)

    c_code2 = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.8), Inches(5.6), Inches(5.0))
    c_code2.fill.solid()
    c_code2.fill.fore_color.rgb = COLOR_CARD
    c_code2.line.color.rgb = COLOR_PURPLE

    tf_c2 = c_code2.text_frame
    tf_c2.word_wrap = True
    p = tf_c2.paragraphs[0]
    p.text = "2. SSE REAL-TIME PROGRESS LISTENER"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    code2 = """function listenToProgress(onUpdate) {
  const es = new EventSource(
    `${API_BASE}/api/events`
  );

  es.onmessage = (event) => {
    const payload = JSON.parse(event.data);
    console.log("Live update:", payload);
    onUpdate(payload);
  };

  return es;
}"""
    p_code2 = tf_c2.add_paragraph()
    p_code2.text = code2
    p_code2.font.size = Pt(11)
    p_code2.font.color.rgb = COLOR_ACCENT
    p_code2.space_before = Pt(10)

    # ----------------------------------------------------
    # SLIDE 10: Critical VPS Rules & Summary
    # ----------------------------------------------------
    slide10 = prs.slides.add_slide(blank_layout)
    add_header(slide10, "9. Critical VPS Operational Rules & Summary")

    rules = [
        ("Mandatory Headless Mode", "Linux VPS servers run without an X11 graphical display window. Always enforce headless=True when launching LinkedInScraper on VPS."),
        ("Linux Process Cleanup Fix", "In app.py, _kill_playwright_chromium() uses Windows 'wmic' and 'taskkill'. On Linux VPS, use cross-platform psutil or 'pkill -f chromium'."),
        ("Server Hardware Recommendation", "Minimum 2GB RAM + 2GB Swap space required on VPS to support headless Chromium memory allocation smoothly."),
        ("Official Repository", "All code, documentation, Dockerfile, and updates are maintained at https://github.com/BawanthaBeliwaththa/Persona_V3")
    ]
    for idx, (rtitle, rdesc) in enumerate(rules):
        row = idx // 2
        col = idx % 2
        x = Inches(0.8 + col * 5.9)
        y = Inches(1.8 + row * 2.6)

        c = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.6), Inches(2.3))
        c.fill.solid()
        c.fill.fore_color.rgb = COLOR_CARD
        c.line.color.rgb = COLOR_BLUE

        tf = c.text_frame
        tf.word_wrap = True
        p_head = tf.paragraphs[0]
        p_head.text = rtitle
        p_head.font.size = Pt(16)
        p_head.font.bold = True
        p_head.font.color.rgb = COLOR_WHITE

        p_body = tf.add_paragraph()
        p_body.text = rdesc
        p_body.font.size = Pt(12)
        p_body.font.color.rgb = COLOR_MUTED
        p_body.space_before = Pt(8)

    output_path = os.path.abspath("Persona_V3_Complete_Guide.pptx")
    prs.save(output_path)
    print(f"Presentation successfully saved to: {output_path}")

if __name__ == '__main__':
    build_presentation()
